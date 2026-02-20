"""Generate a professional .pptx slide deck for Python Turns 35.

Technical theme (dark background, light text) following slides-creator skill specs.
Output can be opened in Google Slides, PowerPoint or Keynote.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

# ---------------------------------------------------------------------------
# Theme colors (Technical theme from .ai/skills/slides-creator)
# ---------------------------------------------------------------------------
BG_COLOR = RGBColor(0x0D, 0x11, 0x17)       # #0D1117 dark blue-black
TEXT_COLOR = RGBColor(0xC9, 0xD1, 0xD9)      # #C9D1D9 light gray
HEADING_COLOR = RGBColor(0x58, 0xA6, 0xFF)   # #58A6FF light blue
ACCENT_YELLOW = RGBColor(0xFF, 0xD4, 0x3B)   # #FFD43B Python yellow
ACCENT_BLUE = RGBColor(0x30, 0x69, 0x98)     # #306998 Python blue
ACCENT_RED = RGBColor(0xE8, 0x58, 0x4A)      # #E8584A warm red
ACCENT_TEAL = RGBColor(0x4B, 0x8B, 0xBE)     # #4B8BBE lighter blue
MUTED_COLOR = RGBColor(0x8B, 0x94, 0x9E)     # #8B949E muted gray
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# Slide dimensions: 16:9 widescreen
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

output_dir = os.path.dirname(os.path.abspath(__file__))


def set_slide_bg(slide, color=BG_COLOR):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=TEXT_COLOR, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    """Add a text box to a slide and return the text frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=18, color=TEXT_COLOR, bold=False,
                  space_before=Pt(6), space_after=Pt(2), alignment=PP_ALIGN.LEFT,
                  font_name="Calibri"):
    """Append a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    p.space_after = space_after
    p.alignment = alignment
    return p


def add_footer(slide, text="Wallace Espindola | speakerdeck.com/wallacese"):
    """Add a subtle footer to the bottom of a slide."""
    add_textbox(slide, Inches(0.5), Inches(6.9), Inches(12), Inches(0.4),
                text, font_size=10, color=MUTED_COLOR,
                alignment=PP_ALIGN.CENTER)


def add_divider(slide, top, color=ACCENT_BLUE):
    """Add a thin horizontal line as a visual divider."""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        Inches(0.8), top, Inches(11.7), Pt(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


# ===== SLIDE BUILDERS =====

def slide_title(prs):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide)

    add_textbox(slide, Inches(1), Inches(1.2), Inches(11), Inches(1.2),
                "Python Turns 35", font_size=48, color=ACCENT_YELLOW, bold=True)

    add_textbox(slide, Inches(1), Inches(2.5), Inches(11), Inches(0.8),
                "The Language That Refused to Die", font_size=28,
                color=WHITE, bold=False)

    add_textbox(slide, Inches(1), Inches(3.4), Inches(11), Inches(0.5),
                "1991  --  2026", font_size=20, color=ACCENT_TEAL)

    add_divider(slide, Inches(4.2))

    add_textbox(slide, Inches(1), Inches(4.6), Inches(11), Inches(0.5),
                "Wallace Espindola", font_size=22, color=WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.5),
                "Senior Software Engineer | Brussels, Belgium",
                font_size=16, color=MUTED_COLOR)
    add_textbox(slide, Inches(1), Inches(5.7), Inches(11), Inches(0.5),
                "speakerdeck.com/wallacese", font_size=14, color=ACCENT_TEAL)


def slide_paradigm_shifts(prs):
    """Slide 2: Few languages survive this long."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "Few Languages Survive This Long", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1),
                     "Most programming languages fade within a decade.",
                     font_size=20, color=WHITE)
    add_paragraph(tf, "Python went through five paradigm shifts and came out "
                  "stronger after each one:", font_size=20, color=TEXT_COLOR,
                  space_before=Pt(16))

    shifts = [
        "Desktop computing",
        "The internet era",
        "Cloud platforms",
        "Mobile",
        "Artificial intelligence",
    ]
    tf2 = add_textbox(slide, Inches(1.2), Inches(3.5), Inches(10), Inches(3),
                      "", font_size=20, color=TEXT_COLOR)
    for i, s in enumerate(shifts):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"  {s}"
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT_YELLOW
        p.font.name = "Calibri"
        p.space_before = Pt(10)

    add_footer(slide)


def slide_phase(prs, year, title, identity, bullets, color):
    """Generic phase slide builder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                f"{title} ({year})", font_size=34,
                color=color, bold=True)

    add_divider(slide, Inches(1.4), color=color)

    tf = add_textbox(slide, Inches(1.0), Inches(1.8), Inches(10.5), Inches(3.5),
                     "", font_size=20, color=TEXT_COLOR)
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"  {b}"
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = "Calibri"
        p.space_before = Pt(12)

    # Identity box
    add_textbox(slide, Inches(0.8), Inches(5.4), Inches(11), Inches(0.6),
                f'Identity: "{identity}"', font_size=22,
                color=color, bold=True)

    add_footer(slide)


def slide_phases(prs):
    """Slides 3-6: Python's four phases."""
    slide_phase(prs,
                "1991-2000", "Phase 1: Scripting & Academia",
                "A better scripting language",
                ["Guido van Rossum created Python as a cleaner alternative to C",
                 "Highly readable, indentation-based syntax",
                 "Built-in standard library for immediate use",
                 "Adopted as a teaching language in universities"],
                ACCENT_BLUE)

    slide_phase(prs,
                "2000-2010", "Phase 2: Web & Open Source",
                "Web productivity language",
                ["Django framework (2005) and Flask (2010)",
                 "Linux ecosystem integration",
                 "Startup adoption wave",
                 "Competing with PHP, Ruby and early Java stacks"],
                ACCENT_TEAL)

    slide_phase(prs,
                "2010-2020", "Phase 3: Data Science & AI",
                "The language of data and AI",
                ["NumPy, Pandas, SciPy for scientific computing",
                 "Jupyter Notebooks changed how we explore data",
                 "TensorFlow and PyTorch standardized ML on Python",
                 "Machine learning became synonymous with Python"],
                ACCENT_YELLOW)

    slide_phase(prs,
                "2020-Present", "Phase 4: Platform & Cloud",
                "Universal integration language",
                ["FastAPI for modern backend services",
                 "Cloud automation and DevOps tooling",
                 "AI orchestration (LangChain, agents, LLMs)",
                 "Cybersecurity tooling and education at global scale"],
                ACCENT_RED)


def slide_survival(prs):
    """Slide 7: Why Python survived."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "Why Did Python Survive 35 Years?", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    factors = [
        ("Human-centered design", "readability over cleverness"),
        ("Batteries included", "standard library for immediate productivity"),
        ("Ecosystem network effects", "universities + startups + researchers"),
        ("Low barrier of entry", "constantly renewing its developer base"),
        ("Adaptability", "absorbed new paradigms instead of resisting them"),
    ]

    tf = add_textbox(slide, Inches(1.0), Inches(1.8), Inches(10.5), Inches(4.5),
                     "", font_size=20, color=TEXT_COLOR)
    for i, (title, desc) in enumerate(factors):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(18)
        run_num = p.add_run()
        run_num.text = f"{i + 1}. "
        run_num.font.size = Pt(20)
        run_num.font.color.rgb = ACCENT_YELLOW
        run_num.font.bold = True
        run_num.font.name = "Calibri"

        run_title = p.add_run()
        run_title.text = title
        run_title.font.size = Pt(20)
        run_title.font.color.rgb = WHITE
        run_title.font.bold = True
        run_title.font.name = "Calibri"

        run_desc = p.add_run()
        run_desc.text = f"  --  {desc}"
        run_desc.font.size = Pt(18)
        run_desc.font.color.rgb = MUTED_COLOR
        run_desc.font.name = "Calibri"

    add_footer(slide)


def slide_strong_points(prs):
    """Slide 8: Python's strong points today."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "Python's Strong Points Today", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    points = [
        ("Developer productivity", "extremely high"),
        ("AI and ML dominance", "the default language for the field"),
        ("Open-source ecosystem", "PyPI with 500k+ packages"),
        ("Cross-domain reach", "web, data, DevOps, security"),
        ("Prototyping speed", "idea to working code, fast"),
        ("Orchestration", "ideal glue language for connecting systems"),
    ]

    tf = add_textbox(slide, Inches(1.0), Inches(1.8), Inches(10.5), Inches(4),
                     "", font_size=20, color=TEXT_COLOR)
    for i, (title, desc) in enumerate(points):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(14)

        run_title = p.add_run()
        run_title.text = f"  {title}"
        run_title.font.size = Pt(20)
        run_title.font.color.rgb = ACCENT_YELLOW
        run_title.font.bold = True
        run_title.font.name = "Calibri"

        run_desc = p.add_run()
        run_desc.text = f"  --  {desc}"
        run_desc.font.size = Pt(18)
        run_desc.font.color.rgb = TEXT_COLOR
        run_desc.font.name = "Calibri"

    add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11), Inches(0.5),
                "Python excels where experimentation matters.",
                font_size=18, color=MUTED_COLOR)

    add_footer(slide)


def slide_opportunities(prs):
    """Slide 9: Opportunities to improve."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "Opportunities to Improve", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    # Left column
    left_items = [
        ("Performance", "No-GIL evolution, better JIT\ncompilation underway"),
        ("Packaging", "pip, poetry, uv, conda --\nstandardization in progress"),
    ]
    # Right column
    right_items = [
        ("Concurrency", "Subinterpreters, experimental\nGIL removal coming"),
        ("Enterprise Governance", "Gradual typing adoption\ngrowing steadily"),
    ]

    y_start = 1.8
    for i, (title, desc) in enumerate(left_items):
        y = y_start + i * 2.2
        add_textbox(slide, Inches(1.0), Inches(y), Inches(5), Inches(0.5),
                    title, font_size=22, color=ACCENT_YELLOW, bold=True)
        add_textbox(slide, Inches(1.0), Inches(y + 0.5), Inches(5), Inches(1.2),
                    desc, font_size=17, color=TEXT_COLOR)

    for i, (title, desc) in enumerate(right_items):
        y = y_start + i * 2.2
        add_textbox(slide, Inches(7.0), Inches(y), Inches(5), Inches(0.5),
                    title, font_size=22, color=ACCENT_YELLOW, bold=True)
        add_textbox(slide, Inches(7.0), Inches(y + 0.5), Inches(5), Inches(1.2),
                    desc, font_size=17, color=TEXT_COLOR)

    add_footer(slide)


def slide_java(prs):
    """Slide 10: Java at 30."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "Meanwhile, Java Turns 30", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(0.8),
                     "Created by James Gosling in 1995 -- another rare survivor.",
                     font_size=20, color=WHITE)
    add_paragraph(tf, "Java solved enterprise-scale problems first:",
                  font_size=20, color=TEXT_COLOR, space_before=Pt(14))

    items = [
        "JVM portability",
        "Strong static typing",
        "Mature concurrency model",
        "Backward compatibility guarantees",
        "Banking, telecom, government systems",
    ]

    tf2 = add_textbox(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(2.5),
                      "", font_size=20, color=TEXT_COLOR)
    for i, item in enumerate(items):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(20)
        p.font.color.rgb = ACCENT_RED
        p.font.name = "Calibri"
        p.space_before = Pt(10)

    add_textbox(slide, Inches(0.8), Inches(5.6), Inches(11), Inches(0.6),
                'Identity: "Engineering-first language."',
                font_size=22, color=ACCENT_RED, bold=True)

    add_footer(slide)


def slide_comparison(prs):
    """Slide 11: Python vs Java comparison table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "Python at 35 vs Java at 30", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    rows = [
        ("Dimension", "Python", "Java"),
        ("Optimization", "Developer speed", "System stability"),
        ("Typing", "Dynamic / Gradual", "Static"),
        ("Performance", "Moderate", "High"),
        ("AI Leadership", "Dominant", "Minimal"),
        ("Concurrency", "Limited historically", "Excellent"),
        ("Learning Curve", "Low", "Moderate"),
        ("Innovation Speed", "Very high", "Controlled"),
    ]

    # Table header and data as textboxes (cleaner than pptx tables)
    col_x = [Inches(1.0), Inches(5.0), Inches(9.0)]
    col_w = [Inches(3.8), Inches(3.8), Inches(3.8)]

    for row_i, (dim, py, ja) in enumerate(rows):
        y = Inches(1.7) + Emu(int(row_i * Inches(0.65)))
        is_header = row_i == 0

        # Alternating row background
        if row_i > 0 and row_i % 2 == 0:
            bg_shape = slide.shapes.add_shape(
                1, Inches(0.8), y - Pt(4),
                Inches(11.7), Inches(0.6)
            )
            bg_shape.fill.solid()
            bg_shape.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
            bg_shape.line.fill.background()

        fs = 16 if is_header else 17
        dim_color = ACCENT_TEAL if is_header else TEXT_COLOR
        py_color = ACCENT_TEAL if is_header else ACCENT_YELLOW
        ja_color = ACCENT_TEAL if is_header else ACCENT_RED

        add_textbox(slide, col_x[0], y, col_w[0], Inches(0.5),
                    dim, font_size=fs, color=dim_color, bold=is_header)
        add_textbox(slide, col_x[1], y, col_w[1], Inches(0.5),
                    py, font_size=fs, color=py_color, bold=is_header)
        add_textbox(slide, col_x[2], y, col_w[2], Inches(0.5),
                    ja, font_size=fs, color=ja_color, bold=is_header)

    add_footer(slide)


def slide_complement(prs):
    """Slide 12: They complement, they don't compete."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "They Don't Compete -- They Complement", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11), Inches(1),
                     "Modern platforms frequently combine both:",
                     font_size=20, color=WHITE)

    add_textbox(slide, Inches(1.2), Inches(2.6), Inches(10), Inches(0.5),
                "Java  for core transactional systems",
                font_size=20, color=ACCENT_RED)
    add_textbox(slide, Inches(1.2), Inches(3.2), Inches(10), Inches(0.5),
                "Python  for AI, analytics, automation and orchestration",
                font_size=20, color=ACCENT_YELLOW)

    add_textbox(slide, Inches(0.8), Inches(4.2), Inches(11), Inches(0.5),
                "Two different optimization layers:",
                font_size=20, color=TEXT_COLOR)

    # Visual blocks
    # Java block
    java_box = slide.shapes.add_shape(
        1, Inches(1.5), Inches(4.9), Inches(4.5), Inches(1.2)
    )
    java_box.fill.solid()
    java_box.fill.fore_color.rgb = RGBColor(0x2A, 0x15, 0x15)
    java_box.line.color.rgb = ACCENT_RED
    java_box.line.width = Pt(1.5)
    add_textbox(slide, Inches(1.7), Inches(5.0), Inches(4), Inches(0.4),
                "Java", font_size=22, color=ACCENT_RED, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.7), Inches(5.4), Inches(4), Inches(0.4),
                "Infrastructure backbone", font_size=16, color=TEXT_COLOR,
                alignment=PP_ALIGN.CENTER)

    # Python block
    py_box = slide.shapes.add_shape(
        1, Inches(7.3), Inches(4.9), Inches(4.5), Inches(1.2)
    )
    py_box.fill.solid()
    py_box.fill.fore_color.rgb = RGBColor(0x1A, 0x1F, 0x15)
    py_box.line.color.rgb = ACCENT_YELLOW
    py_box.line.width = Pt(1.5)
    add_textbox(slide, Inches(7.5), Inches(5.0), Inches(4), Inches(0.4),
                "Python", font_size=22, color=ACCENT_YELLOW, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(7.5), Inches(5.4), Inches(4), Inches(0.4),
                "Intelligence layer", font_size=16, color=TEXT_COLOR,
                alignment=PP_ALIGN.CENTER)

    add_footer(slide)


def slide_future(prs):
    """Slide 13: What's next for both."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "What's Next?", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    # Python column
    add_textbox(slide, Inches(1.0), Inches(1.7), Inches(5), Inches(0.5),
                "Python (Next Decade)", font_size=22,
                color=ACCENT_YELLOW, bold=True)

    py_items = [
        "Gradual typing normalization",
        "Performance evolution (no-GIL, JIT)",
        "AI-native runtime integrations",
        "Stronger packaging standards",
    ]
    tf_py = add_textbox(slide, Inches(1.2), Inches(2.4), Inches(5), Inches(3),
                        "", font_size=18, color=TEXT_COLOR)
    for i, item in enumerate(py_items):
        p = tf_py.paragraphs[0] if i == 0 else tf_py.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = "Calibri"
        p.space_before = Pt(12)

    # Java column
    add_textbox(slide, Inches(7.0), Inches(1.7), Inches(5), Inches(0.5),
                "Java (Next Decade)", font_size=22,
                color=ACCENT_RED, bold=True)

    java_items = [
        "Project Loom concurrency",
        "Cloud-native optimization",
        "Improved developer ergonomics",
        "Continued enterprise dominance",
    ]
    tf_ja = add_textbox(slide, Inches(7.2), Inches(2.4), Inches(5), Inches(3),
                        "", font_size=18, color=TEXT_COLOR)
    for i, item in enumerate(java_items):
        p = tf_ja.paragraphs[0] if i == 0 else tf_ja.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR
        p.font.name = "Calibri"
        p.space_before = Pt(12)

    # Dividing line between columns
    div = slide.shapes.add_shape(
        1, Inches(6.5), Inches(1.7), Pt(1.5), Inches(3.2)
    )
    div.fill.solid()
    div.fill.fore_color.rgb = MUTED_COLOR
    div.line.fill.background()

    add_footer(slide)


def slide_takeaway(prs):
    """Slide 14: Key takeaway."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.9),
                "The Key Takeaway", font_size=36,
                color=HEADING_COLOR, bold=True)

    add_divider(slide, Inches(1.4))

    # Quote block background
    quote_bg = slide.shapes.add_shape(
        1, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.4)
    )
    quote_bg.fill.solid()
    quote_bg.fill.fore_color.rgb = RGBColor(0x16, 0x1B, 0x22)
    quote_bg.line.fill.background()

    # Quote accent bar
    bar = slide.shapes.add_shape(
        1, Inches(1.5), Inches(2.0), Pt(6), Inches(2.4)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT_TEAL
    bar.line.fill.background()

    add_textbox(slide, Inches(2.0), Inches(2.2), Inches(9.5), Inches(0.6),
                'Python survived by optimizing for humans.',
                font_size=24, color=ACCENT_YELLOW, bold=True)

    add_textbox(slide, Inches(2.0), Inches(2.9), Inches(9.5), Inches(0.6),
                'Java survived by optimizing for systems.',
                font_size=24, color=ACCENT_RED, bold=True)

    add_textbox(slide, Inches(2.0), Inches(3.6), Inches(9.5), Inches(0.6),
                "Together -- productivity and reliability.",
                font_size=20, color=WHITE)

    tf = add_textbox(slide, Inches(1.0), Inches(5.0), Inches(11), Inches(1.2),
                     "The future isn't about a single dominant language.",
                     font_size=20, color=TEXT_COLOR)
    add_paragraph(tf, "It's about ecosystems collaborating across strengths.",
                  font_size=20, color=TEXT_COLOR, space_before=Pt(8))

    add_footer(slide)


def slide_closing(prs):
    """Slide 15: Closing / Thank you."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)

    add_textbox(slide, Inches(0.8), Inches(1.0), Inches(11), Inches(1),
                "Thank You", font_size=44, color=ACCENT_YELLOW, bold=True,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6),
                "Happy 35th birthday, Python!",
                font_size=24, color=WHITE, alignment=PP_ALIGN.CENTER)

    add_divider(slide, Inches(2.9))

    add_textbox(slide, Inches(0.8), Inches(3.3), Inches(11), Inches(0.5),
                "Wallace Espindola", font_size=24, color=WHITE, bold=True,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(0.8), Inches(3.9), Inches(11), Inches(0.5),
                "Senior Software Engineer | Brussels, Belgium",
                font_size=16, color=MUTED_COLOR, alignment=PP_ALIGN.CENTER)

    links = [
        ("GitHub", "github.com/wallaceespindola"),
        ("LinkedIn", "linkedin.com/in/wallaceespindola"),
        ("Medium", "medium.com/@wallaceespindola"),
        ("Dev.to", "dev.to/wallaceespindola"),
        ("Speaker Deck", "speakerdeck.com/wallacese"),
    ]

    y = 4.7
    for label, url in links:
        tf = add_textbox(slide, Inches(3.5), Inches(y), Inches(6), Inches(0.35),
                         "", font_size=15, color=TEXT_COLOR,
                         alignment=PP_ALIGN.CENTER)
        p = tf.paragraphs[0]
        run_label = p.add_run()
        run_label.text = f"{label}:  "
        run_label.font.size = Pt(15)
        run_label.font.color.rgb = ACCENT_TEAL
        run_label.font.bold = True
        run_label.font.name = "Calibri"

        run_url = p.add_run()
        run_url.text = url
        run_url.font.size = Pt(15)
        run_url.font.color.rgb = TEXT_COLOR
        run_url.font.name = "Calibri"
        y += 0.38


# ===== MAIN =====

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    slide_title(prs)
    slide_paradigm_shifts(prs)
    slide_phases(prs)          # 4 slides
    slide_survival(prs)
    slide_strong_points(prs)
    slide_opportunities(prs)
    slide_java(prs)
    slide_comparison(prs)
    slide_complement(prs)
    slide_future(prs)
    slide_takeaway(prs)
    slide_closing(prs)

    output_path = os.path.join(output_dir, "python-turns-35-slides.pptx")
    prs.save(output_path)
    print(f"Created: {output_path}")
    print(f"Slides: {len(prs.slides)} total")
    print("Open in Google Slides, PowerPoint or Keynote.")


if __name__ == "__main__":
    main()
